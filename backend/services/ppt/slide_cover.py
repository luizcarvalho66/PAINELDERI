# -*- coding: utf-8 -*-
"""
Slide 1 — Capa executiva estilo Edenred.

Design split-screen com gradiente diagonal, elipse branca,
nome do cliente, período real dos dados e data de geração.
"""

import os
from datetime import datetime
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

from .config import (
    WHITE, EDENRED_RED, EDENRED_RED_DARK, EDENRED_CORAL,
    EDENRED_DARK, TEXT_MUTED, BG_LIGHT, BORDER_GRAY,
    SLIDE_WIDTH, LOGO_PATH, ICONS_DIR, COVER_PATH, MINILOGO_PATH,
    SLIDE_HEIGHT
)
from .helpers import add_text_box


def build_slide_cover(prs, client_name=None, period_start=None, period_end=None):
    """Capa executiva moderna com nome do cliente, período real e data de geração."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ============================================
    # FUNDO GRADIENTE VIBRANTE
    # ============================================
    bg_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg_rect.line.fill.background()
    
    fill = bg_rect.fill
    fill.gradient()
    
    fill.gradient_stops[0].color.rgb = RGBColor(0xE2, 0x06, 0x13)  # Edenred Red
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = RGBColor(0xFF, 0x14, 0x93)  # Magenta/Rosa
    fill.gradient_stops[1].position = 1.0
    
    # Ângulo 135° via XML
    spPr = bg_rect.element.spPr
    gradFill = spPr.find(qn('a:gradFill'))
    if gradFill is not None:
        for old_dir in gradFill.findall(qn('a:lin')) + gradFill.findall(qn('a:path')):
            gradFill.remove(old_dir)
        lin = parse_xml('<a:lin xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" ang="8100000" scaled="1"/>')
        gradFill.append(lin)

    # ============================================
    # SHAPE BRANCO (ELIPSE VERTICAL) DIREITA
    # ============================================
    shape_width = Inches(7.5)  
    shape_height = Inches(9.5) 
    
    center_y = Inches(3.0)
    center_x = SLIDE_WIDTH - Inches(1.5)

    shape_left = center_x - (shape_width / 2)
    shape_top = center_y - (shape_height / 2)

    white_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        shape_left, shape_top,
        shape_width, shape_height
    )
    white_shape.fill.solid()
    white_shape.fill.fore_color.rgb = WHITE
    white_shape.line.fill.background()

    # ============================================
    # LOGO EDENRED DENTRO DO SHAPE BRANCO
    # ============================================
    if os.path.exists(LOGO_PATH):
        # Reduzir o logo e centralizar na parte *visível* do shape branco
        logo_w = Inches(2.8)
        
        # O shape começa em shape_left e vai até SLIDE_WIDTH (na parte visível)
        visible_center_x = shape_left + ((SLIDE_WIDTH - shape_left) / 2)
        
        slide.shapes.add_picture(
            LOGO_PATH,
            visible_center_x - (logo_w / 2), 
            center_y - Inches(0.5), 
            width=logo_w
        )

    # ============================================
    # CONTEÚDO (TEXTOS NA ESQUERDA)
    # ============================================
    text_left = Inches(0.8)
    
    # Título Principal (Linha 1)
    add_text_box(
        slide, text_left, Inches(3.8),
        Inches(10.0), Inches(0.6),
        "Apresentação",
        size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT
    )

    # Título Principal (Linha 2)
    add_text_box(
        slide, text_left, Inches(4.4),
        Inches(10.0), Inches(0.6),
        "Desempenho de RI",
        size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT
    )

    # Linha separadora (Pill curta branca)
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        text_left + Inches(0.05), Inches(5.1), 
        Inches(0.8), Inches(0.06)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = WHITE
    pill.line.fill.background()
    pill.adjustments[0] = 0.5 

    # Nome do Cliente TGM (destaque)
    if client_name:
        add_text_box(
            slide, text_left, Inches(5.3),
            Inches(6.0), Inches(0.5),
            client_name.upper(),
            size=20, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT
        )

    # Período real dos dados
    if period_start and period_end:
        # Formatar datas — aceita string ISO ou objetos date/datetime
        if isinstance(period_start, str):
            try:
                ps = datetime.fromisoformat(period_start).strftime('%d/%m/%Y')
            except ValueError:
                ps = period_start
        else:
            ps = period_start.strftime('%d/%m/%Y')
        
        if isinstance(period_end, str):
            try:
                pe = datetime.fromisoformat(period_end).strftime('%d/%m/%Y')
            except ValueError:
                pe = period_end
        else:
            pe = period_end.strftime('%d/%m/%Y')
        
        periodo_str = f"Período: {ps} — {pe}"
    else:
        periodo_str = ""

    if periodo_str:
        add_text_box(
            slide, text_left, Inches(5.9),
            Inches(6.0), Inches(0.4),
            periodo_str,
            size=14, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT
        )

    # Data de geração do relatório
    data_geracao = datetime.now().strftime('%d/%m/%Y')
    add_text_box(
        slide, text_left, Inches(6.45),
        Inches(6.0), Inches(0.4),
        f"Relatório gerado em {data_geracao}",
        size=12, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT
    )

    # Numeração no canto inferior direito
    add_text_box(
        slide, SLIDE_WIDTH - Inches(0.8), SLIDE_HEIGHT - Inches(0.8),
        Inches(0.5), Inches(0.4),
        "1",
        size=10, bold=False, color=WHITE, alignment=PP_ALIGN.RIGHT
    )

