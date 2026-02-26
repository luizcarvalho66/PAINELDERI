# -*- coding: utf-8 -*-
"""
Slide 2 — KPIs com pills estilo Edenred.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .config import (
    WHITE, EDENRED_RED, TEXT_DARK, TEXT_MUTED, BORDER_GRAY,
    CARD_GREEN, CARD_PURPLE, DANGER_BG, SUCCESS_BG, PURPLE_BG,
)
from .helpers import add_text_box, add_rounded_rect, add_content_header, add_footer_logo, add_pill_badge


def build_slide_kpis(prs, kpi_data: dict):
    """KPIs comprimidos à esquerda (2x3) com espaço livre à direita para análises."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # Novo Header Limpo (Plano Estratégico)
    add_content_header(slide, "Resultados Executivos – RI  |  Últimos 30 Dias")

    import os
    ROOT_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))
    ICONS_KPI_DIR = os.path.join(ROOT_DIR, "assets", "icons_kpi")

    # Atualiza matriz de KPIs para conter o nome do arquivo PNG de ícone exportado
    kpis = [
        ("Análise Total", kpi_data.get("total_analisado", "—"), "Ordens processadas", CARD_PURPLE, PURPLE_BG, "layers.png"),
        ("Economia Real", kpi_data.get("economia_real", "—"), "Valor economizado", CARD_GREEN, SUCCESS_BG, "dollar.png"),
        ("Share Preventiva", kpi_data.get("share_preventiva", "—"), "Mix de manutenções", RGBColor(0x3B, 0x82, 0xF6), RGBColor(0xEF, 0xF6, 0xFF), "pie.png"),
        ("RI Geral", kpi_data.get("ri_geral", "—"), "Média do período", EDENRED_RED, DANGER_BG, "trend.png"),
        ("RI Preventiva", kpi_data.get("ri_preventiva", "—"), "Manutenção programada", CARD_GREEN, SUCCESS_BG, "shield.png"),
        ("RI Corretiva", kpi_data.get("ri_corretiva", "—"), "Manutenção sob demanda", EDENRED_RED, DANGER_BG, "wrench.png"),
    ]

    # Geometria Comprimida (Metade Esquerda)
    card_w = Inches(3.2)
    card_h = Inches(1.5)
    start_x = Inches(0.5)
    start_y = Inches(1.3)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)

    for idx, (title, value, desc, accent, accent_bg, icon_file) in enumerate(kpis):
        # Grade 2 Colunas x 3 Linhas
        col = idx % 2 
        row = idx // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # 1. Base Card - Estética Clean (Apenas Branco e Borda Fina)
        main_card = add_rounded_rect(slide, x, y, card_w, card_h, WHITE, BORDER_GRAY, radius=0.03)

        # 2. Barra Lateral de Acento
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, Inches(0.06), card_h
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()
        accent_bar.adjustments[0] = 0.5 

        # 3. Ícone Profissional PNG 
        icon_path = os.path.join(ICONS_KPI_DIR, icon_file)
        icon_size = Inches(0.4)
        if os.path.exists(icon_path):
            slide.shapes.add_picture(icon_path, x + Inches(0.25), y + Inches(0.20), width=icon_size)

        # 4. Título
        add_text_box(
            slide, x + Inches(0.75), y + Inches(0.25),
            card_w - Inches(0.8), Inches(0.3),
            title.upper(), size=10, bold=True, color=TEXT_MUTED, alignment=PP_ALIGN.LEFT
        )

        # 5. Valor Gigante (Colorido com o acento)
        add_text_box(
            slide, x + Inches(0.25), y + Inches(0.65),
            card_w - Inches(0.5), Inches(0.5),
            str(value), size=28, bold=True, color=accent, alignment=PP_ALIGN.LEFT
        )

        # 6. Descrição Inferior
        add_text_box(
            slide, x + Inches(0.25), y + Inches(1.15),
            card_w - Inches(0.5), Inches(0.3),
            desc, size=10, color=TEXT_MUTED, alignment=PP_ALIGN.LEFT
        )

    # ============================================
    # Adicionar Área do Analista (Metade Direita) - Painel UI
    # ============================================
    right_x = Inches(7.5)
    right_w = Inches(5.3)
    right_h = Inches(5.1)
    right_y = Inches(1.3)
    
    # Fundo do Painel de Síntese
    # Sombra sutil
    add_rounded_rect(slide, right_x+Inches(0.04), right_y+Inches(0.06), right_w, right_h, BORDER_GRAY, line_color=None, radius=0.02)
    # Painel principal
    add_rounded_rect(slide, right_x, right_y, right_w, right_h, WHITE, BORDER_GRAY, radius=0.02)

    # Caixa Título da Análise
    add_text_box(
        slide, right_x, right_y + Inches(0.15),
        right_w, Inches(0.4),
        "SÍNTESE EXECUTIVA", size=14, bold=True, color=TEXT_DARK, alignment=PP_ALIGN.CENTER
    )
    
    # Linha separadora
    add_rounded_rect(slide, right_x + Inches(0.4), right_y + Inches(0.6), right_w - Inches(0.8), Pt(1), BORDER_GRAY, line_color=None, radius=0)

    # Placeholder para digitação
    add_text_box(
        slide, right_x + Inches(0.3), right_y + Inches(0.8),
        right_w - Inches(0.6), right_h - Inches(0.9),
        "• Clique aqui para inserir os seus comentários sobre a evolução dos KPIs...\n"
        "• Detalhe os ofensores do mês e os impactos na operação...\n"
        "• Proponha ações mitigatórias e próximos passos...",
        size=12, color=TEXT_MUTED, alignment=PP_ALIGN.LEFT
    )

    add_footer_logo(slide)

