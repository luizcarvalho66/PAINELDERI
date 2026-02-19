# -*- coding: utf-8 -*-
"""
Slide 6 — Farol de Corretivas — Tabela top chaves.
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .config import (
    WHITE, EDENRED_RED, TEXT_DARK, TEXT_MUTED, LIGHT_GRAY,
    CARD_GREEN, CARD_AMBER, CARD_RED, FONT_FAMILY,
)
from .helpers import add_text_box, add_gradient_header, add_footer_logo


def build_slide_farol_table(prs, farol_data: list):
    """Tabela top chaves do farol."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    add_gradient_header(slide, "Farol de Corretivas — Top Chaves")

    if not farol_data:
        add_text_box(
            slide, Inches(3), Inches(3.5), Inches(7), Inches(1),
            "Sem dados disponíveis para o farol.",
            size=18, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER
        )
        add_footer_logo(slide)
        return

    headers = ["Chave (Peça + MO)", "Qtd OSs", "% Aprovação", "P70 (R$)", "Farol"]
    col_widths = [Inches(4.5), Inches(1.5), Inches(2.0), Inches(2.0), Inches(1.8)]

    n_rows = min(len(farol_data), 15) + 1
    table_top = Inches(1.2)
    table_left = Inches(0.7)
    table_width = sum(w for w in col_widths)
    row_height = Inches(0.38)

    table_shape = slide.shapes.add_table(
        n_rows, len(headers),
        table_left, table_top,
        table_width, row_height * n_rows
    )
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header vermelho
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = EDENRED_RED
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    farol_color_map = {
        "verde": CARD_GREEN,
        "amarelo": CARD_AMBER,
        "vermelho": CARD_RED,
    }

    for row_idx, item in enumerate(farol_data[:15]):
        r = row_idx + 1
        chave = item.get("chave", "—")
        qtd = str(item.get("qtd_os", 0))
        pct = f"{item.get('pct_aprovacao', 0):.1f}%"
        p70 = f"R$ {item.get('p70', 0):,.2f}"
        farol_cor = item.get("farol_cor", "—")
        farol_label = farol_cor.capitalize() if isinstance(farol_cor, str) else "—"

        values = [chave, qtd, pct, p70, farol_label]

        for col_idx, val in enumerate(values):
            cell = table.cell(r, col_idx)
            cell.text = str(val)

            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.name = FONT_FAMILY
                p.font.color.rgb = TEXT_DARK
                p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if col_idx == 4 and isinstance(farol_cor, str):
                fc = farol_color_map.get(farol_cor.lower())
                if fc:
                    for p in cell.text_frame.paragraphs:
                        p.font.bold = True
                        p.font.color.rgb = fc

    add_footer_logo(slide)
