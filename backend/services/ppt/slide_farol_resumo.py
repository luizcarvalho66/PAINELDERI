# -*- coding: utf-8 -*-
"""
Slide 5 — Farol de Corretivas — Resumo & Insights (McKinsey Consulting Style).

Layout:
  - Header: Action Title (vermelho) + sublinhado pill
  - Esquerda: Funnel Infográfico (shapes nativos PowerPoint - vetorial)
  - Direita: Insights analíticos textuais organizados por categoria (Verde/Amarelo/Vermelho)
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from .config import (
    WHITE, TEXT_DARK, FONT_FAMILY,
    CARD_GREEN, CARD_AMBER, CARD_RED,
    ICONS_KPI_DIR,
)
from .helpers import add_text_box, add_content_header, add_footer_logo


# ======================================================================
# CORES E LABELS DO FUNIL
# ======================================================================
_FUNNEL_COLORS = [
    RGBColor(0x47, 0x55, 0x69),  # Slate (Total)
    RGBColor(0x10, 0xB9, 0x81),  # Green (Performando Bem)
    RGBColor(0xF5, 0x9E, 0x0B),  # Amber (Gargalos)
    RGBColor(0xE2, 0x06, 0x13),  # Red (Ação Prioritária)
]

_FUNNEL_LABELS = [
    "TOTAL DE CHAVES",
    "PERFORMANDO BEM",
    "GARGALOS & ATENÇÃO",
    "AÇÃO PRIORITÁRIA",
]


def _add_shadow(shape):
    """Adiciona sombra sutil (outer shadow) via XML direto no shape."""
    spPr = shape._element.spPr
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)

    effectLst = spPr.makeelement(qn('a:effectLst'), {})
    outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
        'blurRad': '63500',
        'dist': '25400',
        'dir': '5400000',
        'algn': 'bl',
    })
    srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = srgbClr.makeelement(qn('a:alpha'), {'val': '35000'})
    srgbClr.append(alpha)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)
    spPr.append(effectLst)


def _add_funnel_native(slide, stats: dict):
    """
    Desenha funil infográfico como shapes NATIVOS do PowerPoint.

    Cada nível é um trapézio vetorial desenhado via FreeformBuilder,
    com preenchimento sólido, sombra sutil e labels adaptativos.

    REGRA: Segmentos são ORDENADOS POR VALOR DECRESCENTE após o Total,
    garantindo que o funil SEMPRE afunila de cima para baixo,
    independente da distribuição dos dados.

    Usa escala sqrt para evitar que valores pequenos gerem
    trapézios estreitos demais para o texto.
    """
    import math

    total_val = stats.get("total", {}).get("value", 0) or 1

    # ---- Montar segmentos e ORDENAR por valor decrescente ----
    segments = [
        {"key": "verde",    "value": stats.get("verde", {}).get("value", 0)},
        {"key": "amarelo",  "value": stats.get("amarelo", {}).get("value", 0)},
        {"key": "vermelho", "value": stats.get("vermelho", {}).get("value", 0)},
    ]
    segments.sort(key=lambda s: s["value"], reverse=True)

    # Mapas de cor e label por key
    _color_map = {
        "verde":    RGBColor(0x10, 0xB9, 0x81),
        "amarelo":  RGBColor(0xF5, 0x9E, 0x0B),
        "vermelho": RGBColor(0xE2, 0x06, 0x13),
    }
    _label_map = {
        "verde":    "PERFORMANDO BEM",
        "amarelo":  "GARGALOS & ATENÇÃO",
        "vermelho": "AÇÃO PRIORITÁRIA",
    }

    # Montar listas ordenadas: Total fixo no topo + segmentos ordenados
    values = [total_val] + [s["value"] for s in segments]
    colors = [_FUNNEL_COLORS[0]] + [_color_map[s["key"]] for s in segments]
    labels = [_FUNNEL_LABELS[0]] + [_label_map[s["key"]] for s in segments]

    # ---- Geometria do funil ----
    center_x = Inches(2.7)
    y_start = Inches(1.55)
    total_h = Inches(5.3)
    gap = Inches(0.06)
    n = 4
    block_h = (total_h - (n - 1) * gap) / n
    max_half_w = Inches(2.5)
    min_half_w = Inches(0.7)  # mínimo generoso p/ legibilidade

    # ---- Escala sqrt: impede estreitamento extremo ----
    def _scaled_hw(val):
        ratio = math.sqrt(val / total_val) if total_val > 0 else 0.2
        return max(min_half_w, int(max_half_w * ratio))

    # Meia-larguras: top de cada trapézio = bottom do anterior
    half_widths = [_scaled_hw(v) for v in values]
    # Base final do último trapézio (mais estreita)
    final_hw = max(min_half_w, int(max_half_w * math.sqrt(values[-1] / total_val) * 0.65))
    half_widths.append(final_hw)

    # Limite de largura abaixo do qual o label fica compacto (só número)
    compact_threshold = Inches(1.3)

    for i in range(n):
        y_top = y_start + i * (block_h + gap)
        y_bot = y_top + block_h

        hw_top = half_widths[i]
        hw_bot = half_widths[i + 1]

        # Vértices do trapézio (sentido horário)
        x_tl = center_x - hw_top
        x_tr = center_x + hw_top
        x_br = center_x + hw_bot
        x_bl = center_x - hw_bot

        # Desenhar via FreeformBuilder
        builder = slide.shapes.build_freeform(x_tl, y_top)
        builder.add_line_segments([
            (x_tr, y_top),
            (x_br, y_bot),
            (x_bl, y_bot),
        ])
        shape = builder.convert_to_shape()

        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i]
        shape.line.fill.background()

        _add_shadow(shape)

        # ---- Textos adaptativos ----
        y_center = (y_top + y_bot) / 2
        perc = int(round(values[i] / total_val * 100))
        avg_hw = (hw_top + hw_bot) / 2
        is_compact = avg_hw < compact_threshold

        if is_compact:
            # COMPACTO: só o número + % (uma linha, fonte menor)
            lbl = slide.shapes.add_textbox(
                center_x - avg_hw + Inches(0.1),
                Emu(int(y_center)) - Inches(0.12),
                avg_hw * 2 - Inches(0.2),
                Inches(0.3),
            )
            p = lbl.text_frame.paragraphs[0]
            p.text = f"{values[i]}  ({perc}%)"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
        else:
            # COMPLETO: nome da categoria + valor abaixo
            lbl = slide.shapes.add_textbox(
                center_x - avg_hw + Inches(0.15),
                Emu(int(y_center)) - Inches(0.22),
                avg_hw * 2 - Inches(0.3),
                Inches(0.3),
            )
            p = lbl.text_frame.paragraphs[0]
            p.text = labels[i]
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER

            sub = slide.shapes.add_textbox(
                center_x - avg_hw + Inches(0.15),
                Emu(int(y_center)) + Inches(0.08),
                avg_hw * 2 - Inches(0.3),
                Inches(0.25),
            )
            p2 = sub.text_frame.paragraphs[0]
            p2.text = f"{values[i]} veículos  ({perc}%)"
            p2.font.size = Pt(10)
            p2.font.color.rgb = WHITE
            p2.font.name = FONT_FAMILY
            p2.alignment = PP_ALIGN.CENTER


def build_slide_farol_resumo(prs, stats: dict, insights: dict | None = None,
                              farol_table_data: list | None = None):
    """
    Constrói o Slide 5 no padrão McKinsey Consulting.

    Args:
        prs: Presentation object
        stats: Dict com chaves verde/amarelo/vermelho/total, cada qual {value: int}
        insights: Dict opcional com chaves verde/amarelo/vermelho, cada qual lista de strings.
                  Se omitido, textos genéricos são usados como placeholder.
        farol_table_data: Lista de dicts do farol (chave, qtd_os, pct_aprovacao, p70, farol_cor).
                          Usado para exibir top exemplos reais em cada card.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    total = stats.get("total", {}).get("value", 0)

    # ==================================================================
    # BACKGROUND BRANDING EDENRED (watermarks suaves — padrão slide 4)
    # ==================================================================

    # Elipse superior direita (rosa claro, suave)
    watermark_top = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9.5), Inches(-2.0),
        Inches(6.0), Inches(6.0),
    )
    watermark_top.fill.solid()
    watermark_top.fill.fore_color.rgb = RGBColor(0xFF, 0xEA, 0xF0)
    watermark_top.line.fill.background()

    # Elipse inferior esquerda (rosa claro, suave)
    watermark_bottom = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-2.0), Inches(4.5),
        Inches(5.0), Inches(5.0),
    )
    watermark_bottom.fill.solid()
    watermark_bottom.fill.fore_color.rgb = RGBColor(0xFF, 0xEA, 0xF0)
    watermark_bottom.line.fill.background()

    # ==================================================================
    # HEADER
    # ==================================================================
    add_content_header(slide, "FAROL DE RI — RESUMO E INSIGHTS")

    # ==================================================================
    # ESQUERDA: FUNNEL NATIVO (shapes vetoriais PowerPoint)
    # ==================================================================
    _add_funnel_native(slide, stats)

    # ==================================================================
    # DIREITA: INSIGHTS EM CARDS PREMIUM + EXEMPLOS REAIS
    # ==================================================================
    default_insights = {
        "verde": [
            "Alta eficiência contínua, volumetria focada na frota leve padrão.",
            "Custos de desgaste natural rodando dentro da curva esperada.",
        ],
        "amarelo": [
            "Incremento nas manutenções de motor em praças isoladas.",
            "Sintoma de padrão de condução inadequado identificado.",
        ],
        "vermelho": [
            "Risco de imobilização na frota pesada das filiais críticas.",
            "Intervenção mecânica ou renegociação de tarifa urgentes.",
        ],
    }
    if insights is None:
        insights = default_insights

    # ---- Separar exemplos reais por cor ----
    exemplos_por_cor = {"verde": [], "amarelo": [], "vermelho": []}
    if farol_table_data:
        for item in farol_table_data:
            cor = (item.get("farol_cor", "") or "").lower()
            if cor in exemplos_por_cor:
                exemplos_por_cor[cor].append(item)
        # Limitar a 2 exemplos por cor (top por qtd_os)
        for cor in exemplos_por_cor:
            exemplos_por_cor[cor] = sorted(
                exemplos_por_cor[cor],
                key=lambda x: x.get("qtd_os", 0), reverse=True,
            )[:2]

    icon_map = {
        "verde": "shield.png",
        "amarelo": "caution.png",
        "vermelho": "alert.png",
    }

    sections = [
        ("PERFORMANDO BEM", "verde", CARD_GREEN, insights.get("verde", [])),
        ("GARGALOS & ATENÇÃO", "amarelo", CARD_AMBER, insights.get("amarelo", [])),
        ("AÇÃO PRIORITÁRIA", "vermelho", CARD_RED, insights.get("vermelho", [])),
    ]

    # Geometria dos cards
    card_x = Inches(5.8)
    card_w = Inches(7.0)
    card_h = Inches(1.85)      # Aumentado para acomodar exemplos reais
    card_y_start = Inches(1.3)
    card_gap = Inches(0.10)
    accent_w = Inches(0.06)
    icon_size = Inches(0.28)

    import os

    from .config import BORDER_GRAY, TEXT_MUTED, LIGHT_GRAY

    for i, (cat_name, cat_key, color_rgb, text_lines) in enumerate(sections):
        curr_y = card_y_start + i * (card_h + card_gap)
        val = stats.get(cat_key, {}).get("value", 0)
        perc = int(round((val / total) * 100)) if total > 0 else 0

        # ---- Container branco arredondado ----
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            card_x, curr_y, card_w, card_h,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER_GRAY
        card.line.width = Pt(1)
        card.adjustments[0] = 0.03

        # ---- Barra accent lateral ----
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            card_x, curr_y, accent_w, card_h,
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = color_rgb
        accent.line.fill.background()

        # ---- Ícone ----
        icon_path = os.path.join(ICONS_KPI_DIR, icon_map.get(cat_key, ""))
        icon_x = card_x + Inches(0.2)
        icon_y = curr_y + Inches(0.12)
        if os.path.exists(icon_path):
            slide.shapes.add_picture(
                icon_path, icon_x, icon_y,
                width=icon_size, height=icon_size,
            )

        # ---- Título da categoria ----
        title_x = icon_x + icon_size + Inches(0.10)
        add_text_box(
            slide, title_x, curr_y + Inches(0.10),
            Inches(5.5), Inches(0.3),
            f"{cat_name} ({perc}%)", size=13, bold=True, color=TEXT_DARK,
        )

        # ---- Bullets analíticos ----
        txt = slide.shapes.add_textbox(
            icon_x + Inches(0.05), curr_y + Inches(0.42),
            card_w - Inches(0.5), Inches(0.45),
        )
        tf = txt.text_frame
        tf.word_wrap = True
        for p_idx, line_str in enumerate(text_lines[:2]):
            p = tf.add_paragraph() if p_idx > 0 else tf.paragraphs[0]
            p.text = f"•  {line_str}"
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_DARK
            p.font.name = FONT_FAMILY
            p.space_after = Pt(2)

        # ---- EXEMPLOS REAIS DE CHAVES ----
        exemplos = exemplos_por_cor.get(cat_key, [])
        if exemplos:
            # Separador fino
            sep_y = curr_y + Inches(0.92)
            sep = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                card_x + Inches(0.2), sep_y,
                card_w - Inches(0.4), Inches(0.01),
            )
            sep.fill.solid()
            sep.fill.fore_color.rgb = BORDER_GRAY
            sep.line.fill.background()

            # Sub-título "Exemplos Reais"
            add_text_box(
                slide, icon_x + Inches(0.05), sep_y + Inches(0.04),
                Inches(3), Inches(0.22),
                "EXEMPLOS REAIS:", size=8, bold=True, color=TEXT_MUTED,
            )

            # Mini-linha para cada chave de exemplo
            ex_y = sep_y + Inches(0.26)
            for j, ex in enumerate(exemplos):
                chave = ex.get("chave", "—")
                qtd = ex.get("qtd_os", 0)
                pct_aprov = ex.get("pct_aprovacao", 0)
                p70 = ex.get("p70", 0)

                # Chip colorido pequeno
                chip = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    icon_x + Inches(0.05), ex_y + Inches(0.04),
                    Inches(0.12), Inches(0.12),
                )
                chip.fill.solid()
                chip.fill.fore_color.rgb = color_rgb
                chip.line.fill.background()
                chip.adjustments[0] = 0.5
                
                # Texto da chave com métricas
                ex_text = f"{chave}   |   {qtd} OSs   |   Aprov: {pct_aprov:.0f}%   |   P70: R$ {p70:,.0f}"
                add_text_box(
                    slide,
                    icon_x + Inches(0.25), ex_y,
                    card_w - Inches(0.7), Inches(0.22),
                    ex_text, size=9, bold=False, color=TEXT_DARK,
                )
                ex_y += Inches(0.24)

    add_footer_logo(slide)
