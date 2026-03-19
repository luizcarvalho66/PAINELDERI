# -*- coding: utf-8 -*-
"""
Slides 3-4 — Charts Plotly como imagem + helper de conversão.
"""

import io
import os
import tempfile
import subprocess

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from .config import WHITE, BORDER_GRAY, FONT_FAMILY, TEXT_MUTED
from .helpers import add_rounded_rect, add_content_header, add_footer_logo

def _prepare_fig_for_ppt(fig_original, is_so_mode=False):
    """Transforma figure interativa do browser em versão PREMIUM para PPT.

    O chart do dashboard usa hoverinfo='none' + tooltip JS externo.
    No PPT (imagem estática), aplicamos um design clean e de alto impacto:
    - Data labels das barras (R$ ou OS) formatados automaticamente
    - Data labels da linha (RI/SO%) em bold e maior destaque
    - Marcadores da linha com borda outline branca (destaca do grid e das barras)
    - Eixos e grids minimalistas (menos poluição visual)

    Args:
        is_so_mode: Se True, formata barras como "XX OS" (Silent Order).
                    Se False, formata como "R$ XXk/M" (monetário).
    """
    import copy
    import plotly.graph_objects as go

    # Clonar a figure para não alterar a original do dashboard
    fig_ppt = go.Figure(fig_original)
    
    # Cores Premium Edenred
    _EDENRED_RED = '#E20613'
    _TEXT_DARK = '#0F172A'
    _TEXT_MUTED = '#64748b'

    for trace in fig_ppt.data:
        if isinstance(trace, go.Bar):
            # Data labels nas barras — usa is_so_mode explícito (não heurística)
            y_vals = trace.y if trace.y is not None else []
            is_os_count = is_so_mode

            text_labels = []
            for v in y_vals:
                try:
                    val = float(v)
                    if is_os_count:
                        text_labels.append(f"<b>{int(val)} OS</b>")
                    elif val >= 1_000_000:
                        text_labels.append(f"<b>R$ {val/1e6:,.1f}M</b>")
                    elif val >= 1_000:
                        text_labels.append(f"<b>R$ {val/1e3:,.0f}k</b>")
                    else:
                        text_labels.append(f"<b>R$ {val:,.0f}</b>")
                except (TypeError, ValueError):
                    text_labels.append("")

            trace.update(
                text=text_labels,
                textposition='inside',      # INSIDE p/ não sobrepor a linha do RI
                insidetextanchor='end',     # Alinha no topo da barra (por dentro)
                marker=dict(
                    color='rgba(226, 6, 19, 0.08)',              # Vermelho pastel muito suave
                    line=dict(color='rgba(226, 6, 19, 0.3)', width=1), # Borda vermelha leve
                ),
                textfont=dict(
                    family=f"{FONT_FAMILY}, sans-serif",
                    size=14,
                    color=_TEXT_DARK,
                ),
                hoverinfo='skip',
                cliponaxis=False,
            )

        elif isinstance(trace, go.Scatter):
            # Data labels nos pontos da linha (RI%)
            y_vals = trace.y if trace.y is not None else []
            text_labels = []
            for v in y_vals:
                try:
                    text_labels.append(f"<b>{float(v):.1f}%</b>")
                except (TypeError, ValueError):
                    text_labels.append("")

            # Preservar cor original ou definir conforme o trace.name (Comparativo vs Geral)
            trace_color = getattr(trace.line, "color", _EDENRED_RED) if trace.line else _EDENRED_RED
            if getattr(trace, "name", "") == "Preventiva":
                trace_color = '#94a3b8' # GREY_LINE do dashboard
            elif getattr(trace, "name", "") == "Corretiva":
                trace_color = _EDENRED_RED

            trace.update(
                mode='lines+markers+text',
                text=text_labels,
                marker=dict(
                    size=12,
                    color=trace_color,
                    symbol='circle',
                    line=dict(color='white', width=2) # Borda branca destaca o ponto
                ),
                line=dict(color=trace_color, width=4, shape='spline'), # Linha mais marcante
                textfont=dict(
                    family=f"{FONT_FAMILY}, sans-serif",
                    size=16,
                    color=trace_color,
                ),
                hoverinfo='skip',
                cliponaxis=False,
            )

    # Lógica Anti-Colisão Dinâmica: se houver 2 traces Scatter, 
    # projeta o texto do marcador mais alto para CIMA e o do mais baixo para BAIXO.
    scatters = [t for t in fig_ppt.data if isinstance(t, go.Scatter)]
    if len(scatters) == 2:
        pos0, pos1 = [], []
        y0_list = scatters[0].y if scatters[0].y is not None else []
        y1_list = scatters[1].y if scatters[1].y is not None else []
        for y0, y1 in zip(y0_list, y1_list):
            v0 = float(y0) if y0 is not None else 0
            v1 = float(y1) if y1 is not None else 0
            if v0 >= v1:
                pos0.append('top center')
                pos1.append('bottom center')
            else:
                pos0.append('bottom center')
                pos1.append('top center')
        scatters[0].update(textposition=pos0)
        scatters[1].update(textposition=pos1)
    else:
        for t in scatters:
            t.update(textposition='top center')

    # Layout otimizado e minimalista para PPT (estático)
    fig_ppt.update_layout(
        font=dict(family=f"{FONT_FAMILY}, sans-serif", size=14, color=_TEXT_DARK),
        title=None, # Título agora é renderizado NATIVAMENTE pelo PPT com Ícone Vetorial
        showlegend=True,
        legend=dict(
            orientation="h",
            y=1.15,
            x=1,
            xanchor="right",
            bgcolor='rgba(255,255,255,0.8)',
            bordercolor='rgba(0,0,0,0.1)',
            borderwidth=1,
            font=dict(size=13, color=_TEXT_DARK),
        ),
        # Eixos mais limpos
        xaxis=dict(
            showgrid=False,
            linecolor='rgba(0,0,0,0.1)',
            linewidth=1,
        ),
        yaxis=dict(
            showgrid=True,
            gridcolor='rgba(0,0,0,0.05)', # Grid bem suave p/ focar no dado
            gridwidth=1,
            zeroline=False,
            visible=False, # Oculta tick labels Y, já temos os valores nos data points
        ),
        yaxis2=dict(
            showgrid=False,
            visible=False, # Oculta tick labels Y2
        ),
        margin=dict(l=40, r=60, t=110, b=90),
        height=560,
        paper_bgcolor='white',
        plot_bgcolor='white',
    )
    
    # Range Y dinâmico: pega o valor máximo real dos dados e adiciona headroom
    try:
        all_y_vals = []
        for trace in fig_ppt.data:
            if hasattr(trace, 'y') and trace.y is not None:
                all_y_vals.extend([float(v) for v in trace.y if v is not None])
        y_max = max(all_y_vals) if all_y_vals else 100
        y_ceil = y_max * 1.25  # 25% headroom para labels não serem cortados
        y_ceil = max(y_ceil, 20)  # Floor mínimo para charts com valores muito baixos
        y_floor = -y_ceil * 0.10  # Desce o eixo Y em -10% visuais para a primeira/última bolha não tocarem no chão
        
        # Como create_comparative_chart não usa make_subplots, secondary_y=False gera erro (silenciado antes).
        # Para ser universal, mexemos no layout root:
        fig_ppt.layout.yaxis.range = [y_floor, y_ceil]
        
        if hasattr(fig_ppt.layout, 'yaxis2'):
            fig_ppt.layout.yaxis2.autorange = True
    except Exception as e:
        print(f"Erro no range Y: {e}")

    return fig_ppt


def fig_to_image_bytes(fig, width=1100, height=520):
    """Converte Plotly Figure para bytes PNG.

    Usa subprocess isolado para evitar conflito asyncio do
    kaleido v1.x + Python 3.14 dentro do Dash.
    """
    fig.update_layout(
        paper_bgcolor='white',
        plot_bgcolor='white',
        font=dict(family=f"{FONT_FAMILY}, sans-serif"),
    )

    fig_json = fig.to_json()
    json_fd, json_path = tempfile.mkstemp(suffix='.json')
    png_fd, png_path = tempfile.mkstemp(suffix='.png')
    os.close(json_fd)
    os.close(png_fd)

    try:
        with open(json_path, 'w', encoding='utf-8') as f:
            f.write(fig_json)

        script = f"""
import plotly.io as pio
import json

with open(r'{json_path}', 'r', encoding='utf-8') as f:
    fig_dict = json.load(f)

fig = pio.from_json(json.dumps(fig_dict))
fig.write_image(r'{png_path}', format='png', width={width}, height={height}, scale=2)
"""
        result = subprocess.run(
            ['python', '-c', script],
            capture_output=True, text=True, timeout=120
        )

        if result.returncode != 0:
            raise RuntimeError(f"Chart export failed: {result.stderr[:500]}")

        with open(png_path, 'rb') as f:
            return f.read()

    finally:
        for p in (json_path, png_path):
            try:
                os.unlink(p)
            except OSError:
                pass


def build_slide_chart(prs, fig, title: str, df_30d=None, silent_orders=None, is_so_mode=False):
    """Slide com chart Plotly como imagem (Design Premium com Respiro e Background Vazado)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # --- Background Abstrato Premium (Watermarks) ---
    # Elipse superior direita (suave, brand rosa)
    watermark_top = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9.5), Inches(-2.0),
        Inches(6.0), Inches(6.0)
    )
    watermark_top.fill.solid()
    watermark_top.fill.fore_color.rgb = RGBColor(0xFF, 0xEA, 0xF0) # Rosa claro (branding)
    watermark_top.line.fill.background()
    
    # Elipse inferior esquerda (suave, brand rosa)
    watermark_bottom = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-2.0), Inches(4.5),
        Inches(5.0), Inches(5.0)
    )
    watermark_bottom.fill.solid()
    watermark_bottom.fill.fore_color.rgb = RGBColor(0xFF, 0xEA, 0xF0)
    watermark_bottom.line.fill.background()
    # ------------------------------------------------

    add_content_header(slide, title)

    # Card branco com borda para o gráfico (altura reduzida para dar respiro ao logo e margem topo aumentada)
    # Antes top=1.2, height=5.3. Agora aumentamos a distância do topo: top=1.7, height=4.8
    add_rounded_rect(
        slide, Inches(0.4), Inches(1.7),
        Inches(7.7), Inches(4.8),
        WHITE, BORDER_GRAY, radius=0.03
    )

    # Preparar figure otimizada para PPT (data labels visíveis)
    fig_ppt = _prepare_fig_for_ppt(fig, is_so_mode=is_so_mode)

    # Chart como imagem (zoom out: mais altura para não clipar labels)
    img_bytes = fig_to_image_bytes(fig_ppt, width=750, height=500)
    img_stream = io.BytesIO(img_bytes)

    # Imagem com padding interno (menor que o card para dar respiro nas bordas)
    slide.shapes.add_picture(
        img_stream,
        Inches(0.5), Inches(1.75),
        width=Inches(7.4), height=Inches(4.7)
    )

    # Adicionar o Painel de Síntese Executiva na Direita
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
    from .config import LIGHT_GRAY, EDENRED_RED, TEXT_DARK, TEXT_MUTED, ICONS_DIR, ICONS_KPI_DIR

    # Fundo do painel de síntese
    # Acompanha a descida do gráfico: top=1.7, height=4.8
    add_rounded_rect(
        slide, Inches(8.4), Inches(1.7),
        Inches(4.4), Inches(4.8),
        LIGHT_GRAY, BORDER_GRAY, radius=0.04
    )
    
    # Decisão de Layout: Mostrar SO Cards se houverem, independentemente do título do chart.
    show_so_cards = bool(silent_orders and len(silent_orders) > 0)

    # Ícone ao lado do título (Síntese Executiva)
    if show_so_cards:
        # Novo ícone gerado pela IA "analysis.png" ou fallback para "warning.png"
        icon_path = os.path.join(ICONS_KPI_DIR, "analysis.png")
        if not os.path.exists(icon_path):
            icon_path = os.path.join(ICONS_KPI_DIR, "warning.png")
            
        if os.path.exists(icon_path):
            slide.shapes.add_picture(icon_path, Inches(8.55), Inches(1.76), width=Inches(0.55))
        title_x = 9.2
        title_w = 3.4
    else:
        icon_path = os.path.join(ICONS_DIR, "executive_summary.png")
        if os.path.exists(icon_path):
            slide.shapes.add_picture(icon_path, Inches(8.55), Inches(1.76), width=Inches(0.55))
            title_x = 9.2
            title_w = 3.4
        else:
            title_x = 8.6
            title_w = 4.0

    # Título do painel (descido para 1.86)
    shape_title = slide.shapes.add_textbox(Inches(title_x), Inches(1.82), Inches(title_w), Inches(0.4))
    tf_title = shape_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "TOP 3 OFENSORES" if show_so_cards else "SÍNTESE EXECUTIVA"
    p_title.font.name = FONT_FAMILY
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = EDENRED_RED
    
    # ------------------------------------------------
    # Título Nativo PPT para o Gráfico com Ícone
    # ------------------------------------------------
    icon_chart_path = os.path.join(ICONS_DIR, "comparative_chart.png")
    if os.path.exists(icon_chart_path):
        slide.shapes.add_picture(icon_chart_path, Inches(0.55), Inches(1.76), width=Inches(0.55))
        chart_title_x = 1.2
    else:
        chart_title_x = 0.55
        
    shape_chart_title = slide.shapes.add_textbox(Inches(chart_title_x), Inches(1.82), Inches(5.0), Inches(0.4))
    tf_chart_title = shape_chart_title.text_frame
    p_chart_title = tf_chart_title.paragraphs[0]
    p_chart_title.text = title.upper() # "Corretiva vs Preventiva" -> CAIXA ALTA
    p_chart_title.font.name = FONT_FAMILY
    p_chart_title.font.size = Pt(13)
    p_chart_title.font.bold = True
    p_chart_title.font.color.rgb = TEXT_DARK

    # Renderizar Cards de Silent Order ou Placeholder
    if show_so_cards:
        from .config import DANGER_BG, WARNING_BG, CARD_RED, CARD_AMBER

        # Subtítulo explicativo
        p_sub = tf_title.add_paragraph()
        p_sub.text = "Estabelecimentos com mais fugas de preventiva."
        p_sub.font.name = FONT_FAMILY
        p_sub.font.size = Pt(8)
        p_sub.font.color.rgb = TEXT_MUTED

        card_top = 2.4 # Top of first card lowered (was 2.0)
        card_colors = [
            (DANGER_BG, CARD_RED),     # 1º pior - vermelho
            (WARNING_BG, CARD_AMBER),  # 2º - amarelo
            (LIGHT_GRAY, TEXT_DARK),   # 3º - neutro
        ]

        for i, so in enumerate(silent_orders[:3]):
            bg_color, accent = card_colors[i] if i < len(card_colors) else card_colors[-1]

            # Card background
            add_rounded_rect(
                slide, Inches(8.55), Inches(card_top),
                Inches(4.1), Inches(1.2),
                bg_color, BORDER_GRAY, radius=0.03
            )

            # Barra lateral de acento
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(8.55), Inches(card_top),
                Inches(0.06), Inches(1.2)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent
            bar.line.fill.background()

            # Nome do estabelecimento
            nome_box = slide.shapes.add_textbox(
                Inches(8.75), Inches(card_top),
                Inches(3.8), Inches(0.35)
            )
            tf_nome = nome_box.text_frame
            tf_nome.word_wrap = True
            p_nome = tf_nome.paragraphs[0]
            p_nome.text = f"{i+1}º {so['nome']}"
            p_nome.font.name = FONT_FAMILY
            p_nome.font.size = Pt(9)
            p_nome.font.bold = True
            p_nome.font.color.rgb = TEXT_DARK

            # Valor principal: % Silent Order
            valor_box = slide.shapes.add_textbox(
                Inches(8.75), Inches(card_top + 0.25),
                Inches(3.8), Inches(0.35)
            )
            tf_val = valor_box.text_frame
            tf_val.word_wrap = True
            p_val = tf_val.paragraphs[0]
            p_val.text = f"{so['so_percent']:.1f}% Fuga"
            p_val.font.name = FONT_FAMILY
            p_val.font.size = Pt(17) # Levemente menor para caber melhor
            p_val.font.bold = True
            p_val.font.color.rgb = accent

            # Linha de detalhe (Total OS | OS sem aprovador)
            det_box = slide.shapes.add_textbox(
                Inches(8.75), Inches(card_top + 0.60),
                Inches(3.8), Inches(0.7)
            )
            tf_det = det_box.text_frame
            tf_det.word_wrap = True
            p_det = tf_det.paragraphs[0]
            p_det.text = f"{so['total_os']} OS no período  ·  {so['so_count']} fugas"
            p_det.font.name = FONT_FAMILY
            p_det.font.size = Pt(8)
            p_det.font.color.rgb = TEXT_DARK

            card_top += 1.3 # Espaçamento entre cards reduzido (era 1.45)
    else:
        # Espaço Livre para Síntese Executiva do Analista
        shape_text = slide.shapes.add_textbox(Inches(8.55), Inches(2.4), Inches(4.1), Inches(4.0))
        tf_text = shape_text.text_frame
        tf_text.word_wrap = True
        
        # Parágrafo 1 - Instrução (mutada, itálico)
        p_text = tf_text.paragraphs[0]
        p_text.text = "Clique aqui para redigir a síntese executiva desta carteira..."
        p_text.font.name = FONT_FAMILY
        p_text.font.size = Pt(10)
        p_text.font.italic = True
        p_text.font.color.rgb = TEXT_MUTED
        
        # Parágrafo 2 - Linha em branco pronta para digitação
        p_blank = tf_text.add_paragraph()
        p_blank.text = ""
        p_blank.font.name = FONT_FAMILY
        p_blank.font.size = Pt(11)
        p_blank.font.italic = False
        p_blank.font.color.rgb = TEXT_DARK

    add_footer_logo(slide)
