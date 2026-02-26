# -*- coding: utf-8 -*-
"""
PPT Helpers — Funções auxiliares compartilhadas entre slides.
"""

import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .config import (
    FONT_FAMILY, TEXT_DARK, WHITE, EDENRED_RED, EDENRED_RED_DARK, EDENRED_CORAL,
    BORDER_GRAY, SLIDE_WIDTH, LOGO_PATH,
)


def set_text(tf, text, size=14, bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT):
    """Define texto num text frame com formatação."""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_FAMILY
    p.alignment = alignment


def add_text_box(slide, left, top, width, height, text, size=14,
                 bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT):
    """Adiciona caixa de texto ao slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    set_text(txBox.text_frame, text, size, bold, color, alignment)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     line_color=None, radius=0.05):
    """Retângulo arredondado."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = radius
    return shape


def add_gradient_header(slide, title_text):
    """(DEPRECATED) Header no estilo Edenred: barra com gradiente vermelho."""
    header_h = Inches(0.85)

    # Base vermelha
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, header_h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = EDENRED_RED
    bar.line.fill.background()

    # Overlay gradiente parcial (canto direito)
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6), Inches(0),
        Inches(7.333), header_h
    )
    fill = overlay.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = EDENRED_RED
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = EDENRED_CORAL
    fill.gradient_stops[1].position = 1.0
    overlay.line.fill.background()

    # Texto do título
    add_text_box(
        slide, Inches(0.8), Inches(0.15),
        Inches(8), Inches(0.6),
        title_text, size=26, bold=True, color=WHITE
    )

def add_content_header(slide, title_text):
    """Header estilo Plano Estratégico (Branco com Título Vermelho e Sublinhado Médio)."""
    # Texto
    add_text_box(
        slide, Inches(0.6), Inches(0.4),
        Inches(10), Inches(0.6),
        title_text.upper(), size=24, bold=True, color=EDENRED_RED
    )
    
    # Sublinhado (Pill Curta Vermelha Subjacente)
    add_rounded_rect(
        slide, Inches(0.65), Inches(0.95), 
        Inches(1.2), Inches(0.08), 
        EDENRED_RED, radius=0.5
    )


def add_footer_logo(slide):
    """Logo Edenred no canto inferior esquerdo."""
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(
            LOGO_PATH,
            Inches(0.4), Inches(6.85),
            width=Inches(0.65)
        )


def add_pill_badge(slide, left, top, width, height, text, bg_color=EDENRED_RED,
                   text_color=WHITE, font_size=12):
    """Pill/badge arredondado no estilo Edenred."""
    pill = add_rounded_rect(slide, left, top, width, height, bg_color, radius=0.5)
    tf = pill.text_frame
    tf.clear()
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = FONT_FAMILY
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    pill.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return pill
